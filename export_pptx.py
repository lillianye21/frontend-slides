import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
# Set 16:9 aspect ratio
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]

BG_COLOR = RGBColor(10, 12, 16)
ACCENT_COLOR = RGBColor(255, 48, 8)
WHITE = RGBColor(255, 255, 255)
GREY = RGBColor(148, 163, 184)

def set_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_text(slide, text, left, top, width, height, font_size, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    return txBox

def add_bullet(slide, title, title_color, body, left, top, width):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    # Title
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.color.rgb = title_color
    p.font.bold = True
    
    # Body
    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.size = Pt(16)
    p2.font.color.rgb = GREY
    
    return txBox

# SLIDE 1: Title
slide1 = prs.slides.add_slide(blank_layout)
set_background(slide1)
# Add Logo
try:
    slide1.shapes.add_picture('DoorDash.png', Inches(1.5), Inches(2.0), height=Inches(0.8))
except Exception as e:
    print("Could not add image:", e)

add_text(slide1, "LILLIAN YE", Inches(1.5), Inches(3.2), Inches(10), Inches(1), 60, ACCENT_COLOR, bold=True)
add_text(slide1, "Candidate for Senior Associate, Marketplace Strategy and Planning", Inches(1.5), Inches(4.3), Inches(10), Inches(0.5), 24, WHITE)
# Red Line
line = slide1.shapes.add_shape(9, Inches(1.5), Inches(5.3), Inches(1.0), Inches(0.02)) # 9 is line
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_COLOR
line.line.color.rgb = ACCENT_COLOR
add_text(slide1, "DOORDASH US MARKETPLACE GROWTH", Inches(1.5), Inches(5.6), Inches(10), Inches(0.5), 14, GREY)

# SLIDE 2: The Truth Seeker
slide2 = prs.slides.add_slide(blank_layout)
set_background(slide2)
add_text(slide2, "Scaling Insights Through Analytical Precision", Inches(1.0), Inches(0.8), Inches(11.333), Inches(1), 36, WHITE, bold=True)

add_bullet(slide2, "Warner Bros. Discovery | Senior Sales Analyst", WHITE, 
           "Streamlined reporting for 80+ international markets, synthesizing complex sales data into actionable insights and dashboards for executive stakeholders.",
           Inches(1.5), Inches(2.0), Inches(10))
add_bullet(slide2, "Team SoloMid | Global Sales Strategy Analyst", WHITE,
           "Used audience and platform data to support partnership planning; helped craft six-figure pitch narratives and materials that brought in new business.",
           Inches(1.5), Inches(3.2), Inches(10))
add_bullet(slide2, "Bank of America | Wealth and Asset Management Analyst", WHITE,
           "Managed high-net-worth portfolios and conducted granular market research for asset allocation across US and International sectors.",
           Inches(1.5), Inches(4.4), Inches(10))
add_bullet(slide2, "Strategic Skill Set", ACCENT_COLOR,
           "Performance Management, Strategic Planning, Data-Informed Storytelling, and Collective Accountability.",
           Inches(1.5), Inches(5.6), Inches(10))

# We need to manually color the specific company names Orange.
# A simpler approach via python-pptx for mixed coloring:
def color_text_run(paragraph, text, color):
    # This is a bit complex in python-pptx, so we will color the titles white, 
    # but the assignment requested companies to be Orange. 
    pass

# To make companies orange, we reconstruct paragraphs
for shape in slide2.shapes:
    if not shape.has_text_frame: continue
    tf = shape.text_frame
    text = tf.paragraphs[0].text
    if " | " in text:
        company, role = text.split(" | ")
        tf.paragraphs[0].text = ""
        run1 = tf.paragraphs[0].add_run()
        run1.text = company + " | "
        run1.font.color.rgb = ACCENT_COLOR
        run1.font.bold = True
        run1.font.size = Pt(20)
        
        run2 = tf.paragraphs[0].add_run()
        run2.text = role
        run2.font.color.rgb = WHITE
        run2.font.bold = True
        run2.font.size = Pt(20)

# SLIDE 3: Operational Excellence
slide3 = prs.slides.add_slide(blank_layout)
set_background(slide3)
add_text(slide3, "Operational Excellence & Process Innovation", Inches(1.0), Inches(0.8), Inches(11.333), Inches(1), 36, WHITE, bold=True)

add_bullet(slide3, "KettleSpace", ACCENT_COLOR, "+37% Profitability\nRefined contract playbooks and workflows, optimizing unit economics across locations.", Inches(1.0), Inches(2.5), Inches(3.5))
tf = slide3.shapes[-1].text_frame
tf.paragraphs[1].font.size = Pt(24)
tf.paragraphs[1].font.bold = True
tf.paragraphs[1].font.color.rgb = WHITE

add_bullet(slide3, "Sonder", ACCENT_COLOR, "15+ Member Teams\nDesigned expansion plans and implemented protocols to enhance guest satisfaction during rapid scale.", Inches(4.9), Inches(2.5), Inches(3.5))
tf = slide3.shapes[-1].text_frame
tf.paragraphs[1].font.size = Pt(24)
tf.paragraphs[1].font.bold = True
tf.paragraphs[1].font.color.rgb = WHITE

add_bullet(slide3, "Core Value", ACCENT_COLOR, "Execute with Precision\nProven ability to bridge the gap between high-level strategy and granular execution.", Inches(8.8), Inches(2.5), Inches(3.5))
tf = slide3.shapes[-1].text_frame
tf.paragraphs[1].font.size = Pt(24)
tf.paragraphs[1].font.bold = True
tf.paragraphs[1].font.color.rgb = WHITE

# SLIDE 4: Tools
slide4 = prs.slides.add_slide(blank_layout)
set_background(slide4)
add_text(slide4, "Optimizing Analytics with Modern Tooling", Inches(1.0), Inches(0.8), Inches(11.333), Inches(1), 36, WHITE, bold=True)

add_bullet(slide4, "The Powerhouse", ACCENT_COLOR, "Leveraging Excel and Tableau as primary drivers for routine analytics and direct data synthesis.", Inches(1.0), Inches(2.5), Inches(3.5))
add_bullet(slide4, "The Edge", ACCENT_COLOR, "Supplementary mastery of Claude and Gemini to automate repetitive tasks and drive operational efficiency.", Inches(4.9), Inches(2.5), Inches(3.5))
add_bullet(slide4, "The Vision", ACCENT_COLOR, "Building robust mechanisms that drive real-time alignment and process improvements across internal stakeholders.", Inches(8.8), Inches(2.5), Inches(3.5))

# SLIDE 5: Strategic Perspective
slide5 = prs.slides.add_slide(blank_layout)
set_background(slide5)
add_text(slide5, "Accelerating the US Marketplace", Inches(1.0), Inches(0.8), Inches(11.333), Inches(1), 36, WHITE, bold=True)

add_bullet(slide5, "Global-to-Local Value", ACCENT_COLOR, "Leveraging international experience (WBD) and US-focused fanbase strategy (TSM) to scale regional partnerships effectively.", Inches(1.5), Inches(2.0), Inches(10))
add_bullet(slide5, "Integrated Execution", ACCENT_COLOR, "Bridging departments to ensure data clarity and process efficiency between regional and international teams.", Inches(1.5), Inches(3.5), Inches(10))
add_bullet(slide5, "Market Precision", ACCENT_COLOR, "Deep dive into current performance trends to find the \"lowest level of detail\" and drive actionable marketplace growth.", Inches(1.5), Inches(5.0), Inches(10))

# SLIDE 6: Roadmap
slide6 = prs.slides.add_slide(blank_layout)
set_background(slide6)
add_text(slide6, "Roadmap for Impact", Inches(1.0), Inches(0.8), Inches(11.333), Inches(1), 36, WHITE, bold=True)

add_bullet(slide6, "30 DAYS", ACCENT_COLOR, "Listen & Learn: Audit internal data systems, build relationships with multidisciplinary teams, and understand performance drivers.", Inches(1.0), Inches(2.5), Inches(3.5))
add_bullet(slide6, "60 DAYS", ACCENT_COLOR, "Optimize & Innovate: Identify analytical bottlenecks; pilot improvements using Tableau/Excel supplemented by AI for speed.", Inches(4.9), Inches(2.5), Inches(3.5))
add_bullet(slide6, "90 DAYS", ACCENT_COLOR, "Execute & Scale: Own performance tracking for major segments; lead the first strategic business planning cycle for Growth.", Inches(8.8), Inches(2.5), Inches(3.5))

# SLIDE 7: Conclusion
slide7 = prs.slides.add_slide(blank_layout)
set_background(slide7)

try:
    slide7.shapes.add_picture('DoorDash.png', Inches(6.0), Inches(0.8), height=Inches(0.8))
except:
    pass

add_text(slide7, "Empowering Local Economies", Inches(1.0), Inches(1.8), Inches(11.333), Inches(1), 36, WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide7, "My mission is to leverage analytical precision and operational grit to accelerate US Marketplace Growth and deliver impact at scale.", Inches(2.0), Inches(2.8), Inches(9.333), Inches(1), 20, GREY, align=PP_ALIGN.CENTER)

# Contact Box
# add a rectangle
shape = slide7.shapes.add_shape(1, Inches(4.0), Inches(4.0), Inches(5.333), Inches(2.5))
shape.fill.background()
shape.line.color.rgb = ACCENT_COLOR

tf = shape.text_frame
tf.word_wrap = True
tf.clear()

p1 = tf.paragraphs[0]
p1.text = "Lillian Ye"
p1.font.size = Pt(24)
p1.font.color.rgb = ACCENT_COLOR
p1.font.bold = True
p1.alignment = PP_ALIGN.LEFT

p2 = tf.add_paragraph()
p2.text = "Candidate for Senior Associate, Marketplace Strategy and Planning"
p2.font.size = Pt(16)
p2.font.color.rgb = WHITE
p2.alignment = PP_ALIGN.LEFT

p3 = tf.add_paragraph()
p3.text = "\nLillian.ye21@gmail.com\nhttps://www.linkedin.com/in/lillian-ye/"
p3.font.size = Pt(14)
p3.font.color.rgb = GREY
p3.alignment = PP_ALIGN.LEFT

prs.save('doordash-interview.pptx')
print("Successfully saved doordash-interview.pptx!")
